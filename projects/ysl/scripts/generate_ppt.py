#!/usr/bin/env python3
"""SKC PPT를 베이스로 YSL 데이터 교체.

전략:
1. SKC PPT 복사 → YSL 파일
2. 전역 텍스트 치환 (브랜드명, 제품명, 카테고리, URL 등)
3. 슬라이드별 데이터 표/차트 교체
4. 데이터 부족 슬라이드는 노란 TBD 박스 오버레이
"""
import json
import re
import shutil
from copy import deepcopy
from pathlib import Path

import pandas as pd
from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

ROOT = Path('/Users/kimhansol/dev/ysl')
SKC_PPT = ROOT / 'BubbleShare_Skinceuticals_GEO Audit Report_1st_0416.pptx'
OUT_PPT = ROOT / 'final' / 'BubbleShare_YSL_GEO_Audit_Report_1st.pptx'
QLIST = ROOT / 'final' / 'Bubbleshare_YSL_Question_List_수정.xlsx'

# 메트릭 로드
metrics = json.loads(Path('/tmp/ysl_metrics.json').read_text(encoding='utf-8'))
mention_df = pd.read_pickle('/tmp/ysl_mention.pkl')
citation_df = pd.read_pickle('/tmp/ysl_citation.pkl')


# Question List에서 추가 메타 (축1/2/3, 인텐트 정의 등) 로드
def load_qlist_extras():
    wb = load_workbook(QLIST, data_only=True)
    ws = wb['Questions (Final)']
    rows = []
    for r in range(5, 365):
        no = ws.cell(r, 2).value
        if isinstance(no, (int, float)):
            rows.append({
                'no': int(no),
                'category': ws.cell(r, 3).value,
                'sub': ws.cell(r, 4).value,
                'intent': ws.cell(r, 5).value,
                'question': ws.cell(r, 6).value,
                'axis1': ws.cell(r, 8).value,
                'axis2': ws.cell(r, 9).value,
                'axis3': ws.cell(r, 10).value,
            })
    qmeta = pd.DataFrame(rows)

    # Prompt setting에서 카테고리 정의·인텐트 정의
    ws_ps = wb['Prompt setting']
    cat_defs = {
        '향수': ws_ps.cell(29, 4).value,
        '기프팅': ws_ps.cell(30, 4).value,
        '쿠션': ws_ps.cell(31, 4).value,
    }
    intent_defs = []
    for r in [37, 38, 39, 40]:
        intent_defs.append({
            'name': ws_ps.cell(r, 2).value,
            'def': ws_ps.cell(r, 3).value,
            'examples': ws_ps.cell(r, 4).value,
            'pattern': ws_ps.cell(r, 5).value,
            'ratio': ws_ps.cell(r, 6).value,
        })

    return qmeta, cat_defs, intent_defs


qmeta_df, CAT_DEFS, INTENT_DEFS = load_qlist_extras()


# =============================================================================
# 전역 텍스트 치환 사전
# =============================================================================
GLOBAL_REPLACEMENTS = [
    # 브랜드명
    ('SkinCeuticals', 'YSL Beauty'),
    ('skinceuticals', 'yslbeauty'),
    ('SKINCEUTICALS', 'YSL BEAUTY'),
    ('스킨수티컬즈', '입생로랑'),
    ('SKC', 'YSL'),
    # 제품/시그니처
    ('CE 페룰릭', 'Libre'),
    ('CE페룰릭', 'Libre'),
    ('CE Ferulic', 'Libre EDP'),
    # SKC 경쟁사 → 빈 문자열 또는 안전 placeholder
    # (이들은 SKC의 경쟁사이고 YSL과 다른 카테고리라 일괄 치환 시 의미가 깨짐)
]

# 슬라이드 prose 내 SKC 특정 키워드 — 이 키워드 포함된 텍스트박스는 TBD로 wholesale 교체
SKC_PROSE_KEYWORDS = [
    # 경쟁사
    'COSRX', 'VT', '닥터디퍼런트', '폴라초이스', '폴라스 초이스',
    'CEQ', '리쥬란', '구달', '피지오겔', '라로슈포제', '메디큐브',
    # SKC 카테고리/콘셉트 (YSL과 무관)
    '시술 후', '시술 전', '시술 대신', '시술 경험', '시술 관련', '시술명',
    '향수 입문자', '항산화 앰플', '비타민C 앰플', '비타민C 입문', '앰플',
    '골든룰', '쉘던 핀넬',
    # SKC 인텐트 라벨 (YSL은 4 인텐트 다름)
    'How-to', 'Concern',
    # SKC 저니/카테고리 (YSL은 3 카테고리: 향수/기프팅/쿠션)
    '피부 고민', '성분/기술력', '경쟁/가격 비교',
]

# SKC 특정 수치 — prose에 이 수치가 포함되면 SKC 잔존으로 판정
SKC_NUMBERS = [
    '5,430', '14,224', '12,053',  # SKC 응답·인용 카운트
    '1.8%', '3.4%', '3.1%', '14.1%', '0.33%', '1.73%',  # SKC 핵심 멘션률
    '99건', '248건', '655', '1,660',  # SKC 카운트
    '76건', '215건', '60건',
]
TBD_PROSE = '[TBD] 인사이트 카피 작성 필요 — 1차 분석 결과 기반으로 수아 검토 후 확정'

# 카테고리/제품 매핑은 신중하게 처리 — 일부만
CATEGORY_REPLACEMENTS = [
    ('항산화 앰플', '럭셔리 향수'),
    ('항산화', '향수'),
]

# 채널 표시 라벨 — 데이터의 lowercase 키와 PPT 표시 라벨 매핑
CHANNEL_LABEL = {
    'chatgpt': 'ChatGPT',
    'google': 'Google AIO',
    'naver': 'Naver',
}

# YSL Beauty 분석 컨텍스트
YSL_CONTEXT = {
    'brand': 'YSL Beauty',
    'brand_kr': '입생로랑 뷰티',
    'analysis_period': '2026년 4월',
    'total_questions': 360,
    'total_responses': metrics['total_responses'],
    'total_citations': metrics['total_citations'],
    'ai_response_rate': metrics['ai_response_rate'],
    'overall_mention_rate': metrics['funnel_all']['brand_mention_rate'],
    'top_competitor': max((b, r) for b, r in metrics['sov_overall'].items() if b != 'YSL Beauty'),
}

# Categories
SKC_TO_YSL_INTENT = {
    # SKC's intent labels likely don't match YSL exactly. Keep as-is for now.
}


# =============================================================================
# 헬퍼
# =============================================================================

def replace_text_in_runs(text_frame, replacements):
    """text_frame 내 모든 run에 대해 치환. 다중 run에 걸친 텍스트도 처리."""
    for para in text_frame.paragraphs:
        # 1. 단일 run 내 치환 (1차)
        for run in para.runs:
            new = run.text
            for old, replacement in replacements:
                new = new.replace(old, replacement)
            if new != run.text:
                run.text = new

        # 2. 다중 run에 걸친 패턴 처리 — 결합 텍스트 검사
        full = ''.join(r.text for r in para.runs)
        replaced = full
        for old, replacement in replacements:
            replaced = replaced.replace(old, replacement)
        if replaced != full and para.runs:
            # 첫 run에 전체 결과 넣고 나머지 비우기
            para.runs[0].text = replaced
            for r in para.runs[1:]:
                r.text = ''


def replace_text_in_shape(shape, replacements):
    if shape.has_text_frame:
        replace_text_in_runs(shape.text_frame, replacements)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                replace_text_in_runs(cell.text_frame, replacements)
    if hasattr(shape, 'shapes'):
        for inner in shape.shapes:
            replace_text_in_shape(inner, replacements)


def scrub_skc_prose(slide):
    """SKC 특정 키워드 또는 수치 포함 텍스트박스/표 셀을 TBD placeholder로 교체."""
    def has_skc_keyword(text):
        if any(kw in text for kw in SKC_PROSE_KEYWORDS):
            return True
        # 숫자는 prose 길이 30자 이상에서만 검출 (data table 셀 false positive 방지)
        if len(text.strip()) > 30 and any(n in text for n in SKC_NUMBERS):
            return True
        return False

    def clear_text_frame(tf, replacement=''):
        for para in tf.paragraphs:
            for run in para.runs:
                run.text = ''
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = replacement
        elif tf.paragraphs:
            run = tf.paragraphs[0].add_run()
            run.text = replacement

    def scrub_shape(shape):
        # 텍스트박스
        if shape.has_text_frame and not shape.has_table:
            full = shape.text_frame.text
            stripped = full.strip()
            if has_skc_keyword(full):
                if len(stripped) > 30:
                    clear_text_frame(shape.text_frame, TBD_PROSE)
                else:
                    # 짧은 라벨 — keyword/숫자만 있는 경우 clear
                    if any(kw in full for kw in SKC_PROSE_KEYWORDS) or any(n in full for n in ['5,430', '14,224', '12,053', '99건', '248건']):
                        clear_text_frame(shape.text_frame, '')
            # 추가: 짧은 라벨도 SKC 절대 수치 포함 시 clear
            elif stripped and any(n in stripped for n in ['5,430', '14,224', '12,053', '655', '1,660', '248건']):
                clear_text_frame(shape.text_frame, '')
        # 표 셀
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    full = cell.text_frame.text
                    if has_skc_keyword(full):
                        if len(full.strip()) > 30:
                            clear_text_frame(cell.text_frame, '[TBD]')
                        else:
                            if any(kw in full for kw in SKC_PROSE_KEYWORDS):
                                clear_text_frame(cell.text_frame, '')
        # Group shape 내부도
        if shape.shape_type == 6:
            try:
                for inner in shape.shapes:
                    scrub_shape(inner)
            except Exception:
                pass

    for shape in slide.shapes:
        scrub_shape(shape)


def add_tbd_overlay(slide, message='[TBD] 1차 데이터 분석 후 수아 리뷰 → 보강 필요', position='top'):
    """슬라이드 상단에 노란 TBD 박스 추가."""
    if position == 'top':
        left, top, width, height = Emu(914400 * 0.5), Emu(914400 * 0.3), Emu(914400 * 12.3), Emu(914400 * 0.4)
    else:
        left, top, width, height = Emu(914400 * 0.5), Emu(914400 * 6.8), Emu(914400 * 12.3), Emu(914400 * 0.5)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
    box.line.color.rgb = RGBColor(0xE6, 0xB8, 0x00)
    box.line.width = Pt(1)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = message
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x80, 0x60, 0x00)


def fill_table(table, data_rows, header_row=True, start_row=1, clear_remaining=True, clear_extra_cols=False):
    """주어진 표를 data_rows로 채우기. data_rows = [[셀1, 셀2, ...], ...]
    clear_remaining=True: data_rows 이후 남은 행은 비우기
    clear_extra_cols=True: 각 행에서 data 컬럼 이후 남은 셀도 비우기
    """
    last_filled = start_row - 1
    for i, row_data in enumerate(data_rows):
        r = start_row + i
        if r >= len(table.rows):
            break
        last_filled = r
        for j, val in enumerate(row_data):
            if j >= len(table.columns):
                break
            cell = table.rows[r].cells[j]
            tf = cell.text_frame
            for para in tf.paragraphs:
                for run in para.runs:
                    run.text = ''
            if tf.paragraphs and tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = str(val)
            else:
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
        # Clear extra columns in this row
        if clear_extra_cols:
            for j in range(len(row_data), len(table.columns)):
                cell = table.rows[r].cells[j]
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''
    # Clear remaining rows (after last filled)
    if clear_remaining:
        for r in range(last_filled + 1, len(table.rows)):
            for j in range(len(table.columns)):
                cell = table.rows[r].cells[j]
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''


# =============================================================================
# 슬라이드별 처리
# =============================================================================

def process_slide(idx, slide, prs):
    """슬라이드 인덱스(1-based)에 따라 처리."""
    # 1. 전역 텍스트 치환 (모든 슬라이드)
    for shape in slide.shapes:
        replace_text_in_shape(shape, GLOBAL_REPLACEMENTS)
        replace_text_in_shape(shape, CATEGORY_REPLACEMENTS)
        # Group shapes 내부도 처리
        if shape.shape_type == 6:  # Group
            try:
                for inner in shape.shapes:
                    replace_text_in_shape(inner, GLOBAL_REPLACEMENTS)
                    replace_text_in_shape(inner, CATEGORY_REPLACEMENTS)
            except Exception:
                pass

    # 1-b. SKC 특정 prose 스크러빙 — 인사이트 텍스트박스 wholesale 교체
    scrub_skc_prose(slide)

    # 2. 슬라이드별 추가 작업
    if idx == 1:
        # Cover — 날짜를 한국어 형식으로
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if 'April 2026' in t:
                    replace_text_in_runs(shape.text_frame, [('April 2026', '2026년 4월')])
                if '1st Analysis report' in t:
                    replace_text_in_runs(shape.text_frame, [
                        ('1st Analysis report', '1차 분석 보고서'),
                        ('Shaping', 'AI 시대의 럭셔리 뷰티 — '),
                        ('leadership at the moment AI influences consumer decisions', '리더십 강화 전략'),
                        ('GEO AI Visibility Audit', 'GEO AI Visibility Audit — 1차 진단'),
                    ])

    elif idx == 4:
        _fill_slide_4_overview(slide)

    elif idx == 5:
        _fill_slide_5_design(slide)

    elif idx == 8:
        # 퍼널 4표: 전체 / 채널 / 카테고리 / 인텐트
        # SKC 표 구조: 7x4 / 7x4 / 7x7 / 7x6
        # YSL 데이터로 가능한 것만 채움
        # (table replacement will be done after global text replacement)
        _fill_slide_8_funnel(slide)
        add_tbd_overlay(slide, '데이터 보강 필요: Perplexity/Gemini/Google AI Mode 등 6채널 가용성 확인 후 채널 컬럼 보강', position='bottom')

    elif idx == 9:
        # 인용 도메인 표 2개
        _fill_slide_9_citation(slide)

    elif idx == 11:
        # 채널 × 경쟁사 Mention Rate (6x19)
        _fill_slide_11_competitor(slide)
        add_tbd_overlay(slide, '데이터 보강: 채널 6개 중 3개만 가용 (ChatGPT/Google/Naver). Perplexity/Gemini 미수집', position='bottom')

    elif idx == 12:
        _fill_slide_12_engine_compare(slide)

    elif idx == 13:
        _fill_slide_13_intent_journey(slide)

    elif idx == 14:
        # 카테고리 딥다이브 (브랜드 랭킹)
        _fill_slide_14_category_deep_dive(slide)

    elif idx == 15:
        # 페르소나/스킨타입 — 축2 기반
        _fill_slide_15_persona(slide)

    elif idx == 16:
        # USP 차별화 — 데이터 부족
        # 모든 표 셀과 텍스트박스를 [TBD]로 강제 클리어
        for shape in slide.shapes:
            if shape.has_text_frame and not shape.has_table:
                t = shape.text_frame.text.strip()
                if len(t) > 30:
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ''
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = TBD_PROSE
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text_frame.text.strip()
                        # 모든 데이터 셀 클리어 (헤더는 두지만 SKC 절대값 포함 시 클리어)
                        if any(kw in t for kw in ['99건', '76건', '215건', '60건', 'Libre vs', '닥터디퍼런트']):
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.text = ''
                        elif len(t) > 5 and 'YSL' not in t and '입생로랑' not in t and t not in ['차이점', '점수', 'AI', '비고', '유무']:
                            for para in cell.text_frame.paragraphs:
                                for run in para.runs:
                                    run.text = ''
        add_tbd_overlay(slide, '데이터 부족: YSL 시그니처 제품(Libre/Y/Mon Paris) 공식 USP 메시지 자료 필요. AI 인지 여부 분석 결과만 자동 추출 가능', position='top')

    elif idx == 17:
        # 포지셔닝 맵 — 표 + 버블 라벨
        _fill_slide_17_positioning(slide)
        _fill_slide_17_bubbles(slide)
        add_tbd_overlay(slide, '버블 차트 위치는 수동 조정 필요 (X: 인지 Mention %, Y: 전환 Mention %)', position='bottom')

    elif idx == 19:
        # SOV 랭킹 3표
        _fill_slide_19_sov(slide)

    elif idx == 20:
        add_tbd_overlay(slide, '대체 컨텍스트 분석은 텍스트 분석 + 수동 검수 필요 — 1차 게재 보류', position='top')

    elif idx == 21:
        # 경쟁사 직접 대체 분석 — SKC 특화 분석. 표 + TBD 텍스트박스 모두 삭제
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.has_table:
                shapes_to_remove.append(shape)
            elif shape.has_text_frame and 'TBD' in shape.text_frame.text and '인사이트' in shape.text_frame.text:
                # 이전 라운드에 추가된 [TBD] 인사이트 텍스트박스 제거
                shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

        # 큰 안내 박스 추가 (슬라이드 중앙)
        from pptx.util import Inches
        msg_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2.5), Inches(11.3), Inches(2.5))
        msg_box.fill.solid()
        msg_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF2, 0xCC)
        msg_box.line.color.rgb = RGBColor(0xE6, 0xB8, 0x00)
        tf = msg_box.text_frame
        tf.word_wrap = True
        from pptx.enum.text import MSO_ANCHOR
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p_msg = tf.paragraphs[0]
        from pptx.enum.text import PP_ALIGN
        p_msg.alignment = PP_ALIGN.CENTER
        run = p_msg.add_run()
        run.text = '[데이터 부족] 경쟁사 직접 대체 분석\n\n"AI가 YSL 대신 X 브랜드를 추천한 케이스" 자동 추출은 가능하나 텍스트 분석 + 수동 검수 필요.\n\n1차 게재 보류 — 2차 분석 단계에서 보강 예정.'
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x80, 0x60, 0x00)

    elif idx == 23:
        # 차트 2개 (도메인 / 콘텐츠 타입)
        _fill_slide_23_citation_charts(slide)

    elif idx == 24:
        # 엔진별 Top 도메인 표 3개
        _fill_slide_24_top_domains(slide)

    elif idx == 26:
        # 이커머스 인용
        _fill_slide_26_ecommerce(slide)

    elif idx == 27:
        # 자사몰 인용
        _fill_slide_27_brand_mall(slide)
        add_tbd_overlay(slide, 'YSL 자사몰 모범 사례 비교 케이스(Polars Choice 같은) 정보 필요', position='bottom')

    elif idx == 28:
        # 레퍼런스 이미지
        add_tbd_overlay(slide, '데이터 부족: YSL 자사몰 모범 사례 스크린샷 또는 비교 대상 케이스 필요', position='top')

    elif idx == 30:
        # 액션 플랜
        add_tbd_overlay(slide, 'YSL 맞춤 액션 플랜 카피라이팅은 수아 리뷰 후 확정 — 1차 게재 보류', position='top')

    elif idx == 32:
        # 분석 각도 (축1/2/3)
        _fill_slide_32_analysis_angle(slide)

    elif idx == 33:
        # 인텐트 정의
        _fill_slide_33_intent_def(slide)

    elif idx == 34:
        # 벤치마크 노트 (KW vs Q)
        _fill_slide_34_benchmark(slide)

    elif idx == 35:
        # 이슈 레지스터 — SKC 이슈 항목 모두 클리어
        for shape in slide.shapes:
            if shape.has_table:
                for r_idx in range(1, len(shape.table.rows)):
                    for c_idx in range(len(shape.table.columns)):
                        cell = shape.table.rows[r_idx].cells[c_idx]
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ''
                # R1C0에 TBD 표시
                if len(shape.table.rows) >= 2:
                    cell = shape.table.rows[1].cells[0]
                    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].text = '[TBD] 25개 이슈 항목 — 1차 분석 결과 종합 후 수동 작성'
                    elif cell.text_frame.paragraphs:
                        cell.text_frame.paragraphs[0].add_run().text = '[TBD] 25개 이슈 항목'
        add_tbd_overlay(slide, '이슈 레지스터는 분석 결과 종합 후 수동 작성 — 1차 게재 보류', position='top')


# =============================================================================
# 슬라이드별 데이터 채우기
# =============================================================================

def _fill_slide_4_overview(slide):
    """슬라이드 4: 분석 개요 (타임라인·플랫폼·쿼리·메트릭)."""
    tables = [s for s in slide.shapes if s.has_table]

    total_q = 360
    total_resp = metrics['total_responses']  # 1080
    total_cit = metrics['total_citations']  # 2839

    # 표 #0 (2x3): 타임라인 / 플랫폼 / 쿼리 — Header + Data row
    if len(tables) >= 1:
        t = tables[0].table
        if len(t.rows) >= 2 and len(t.columns) >= 3:
            data = [
                '분석: 2026-04-24 ~ 04-27\n수집 사이클: 2회',
                'ChatGPT\nGoogle (AI Overview)\nNaver (Search)',
                f'{total_q}건 질문\n+ {total_q}건 키워드',
            ]
            for c_idx, val in enumerate(data):
                cell = t.rows[1].cells[c_idx]
                tf = cell.text_frame
                # Clear all
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ''
                # Set first paragraph first run
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = val
                elif tf.paragraphs:
                    p = tf.paragraphs[0].add_run()
                    p.text = val

    # 표 #11 (8x2): Metric / Value
    # 4번째 표 정도 (일반적으로 마지막 표가 Metric)
    metric_table = None
    for tbl_shape in tables:
        t = tbl_shape.table
        if len(t.rows) >= 8 and len(t.columns) == 2:
            # 첫 셀이 'Metric'
            if 'Metric' in t.rows[0].cells[0].text_frame.text:
                metric_table = t
                break
    if metric_table is None:
        return

    metric_data = [
        ('Unique Queries', f'{total_q}'),
        ('Total Responses', f'{total_resp:,}'),
        ('Citation', f'{total_cit:,}'),
        ('Channels', 'ChatGPT, Google (AIO), Naver (Search)'),
        ('Responses per Query (avg)', f'{total_resp/total_q:.1f}'),
        ('Analysis Period', '2026.04.24 ~ 2026.04.27'),
        ('Brand', 'YSL Beauty (입생로랑)'),
    ]
    for i, (k, v) in enumerate(metric_data):
        r = i + 1  # row 0 is header
        if r >= len(metric_table.rows):
            break
        # Key cell (col 0)
        cell_k = metric_table.rows[r].cells[0]
        for para in cell_k.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''
        if cell_k.text_frame.paragraphs and cell_k.text_frame.paragraphs[0].runs:
            cell_k.text_frame.paragraphs[0].runs[0].text = k
        elif cell_k.text_frame.paragraphs:
            cell_k.text_frame.paragraphs[0].add_run().text = k
        # Value cell (col 1)
        cell_v = metric_table.rows[r].cells[1]
        for para in cell_v.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''
        if cell_v.text_frame.paragraphs and cell_v.text_frame.paragraphs[0].runs:
            cell_v.text_frame.paragraphs[0].runs[0].text = v
        elif cell_v.text_frame.paragraphs:
            cell_v.text_frame.paragraphs[0].add_run().text = v


def _fill_slide_5_design(slide):
    """슬라이드 5: 분석 설계 — YSL 카테고리(향수/기프팅/쿠션) + 4 인텐트 프레임워크."""
    tables = [s for s in slide.shapes if s.has_table]
    # 표 8x4: 카테고리 × 목적 × 비율
    if len(tables) >= 2:
        target = None
        for tbl_shape in tables:
            t = tbl_shape.table
            if len(t.rows) >= 7 and len(t.columns) == 4:
                target = t
                break
        if target:
            data = [
                ['향수 (Perfume)', '럭셔리 향수 카테고리 — 여성/남성/계절·TPO·선물 상황·지속력/확산력',
                 'Libre·Y·Mon Paris·Black Opium 등 시그니처 라인 노출 진단', '40%'],
                ['기프팅 (Gifting)', '뷰티 선물 — 관계별(여친/엄마/시어머니), 예산별(3~20만원), 채널별(카카오/면세점/백화점)',
                 '선물 의사결정에서 YSL 등장 비율, 호불호 없는 선물 포지셔닝', '40%'],
                ['쿠션 (Cushion)', '쿠션 파운데이션 — 피부타입/계절/연령/명품 vs 로드샵 비교',
                 'Touche Éclat Le Cushion 노출 + 명품 쿠션 시장 점유율', '20%'],
            ]
            # 헤더 R0 두고 R1~R3 채우고 R4 이후 클리어
            for i, row_data in enumerate(data):
                r = 1 + i
                if r >= len(target.rows):
                    break
                for c, val in enumerate(row_data):
                    if c >= len(target.columns):
                        break
                    cell = target.rows[r].cells[c]
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''
                    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].text = str(val)
                    elif cell.text_frame.paragraphs:
                        cell.text_frame.paragraphs[0].add_run().text = str(val)
            # 남은 행 클리어 (Total 빼고)
            for r in range(4, len(target.rows) - 1):
                for c in range(len(target.columns)):
                    cell = target.rows[r].cells[c]
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''
            # Total 행 (마지막 행) — 100%
            if len(target.rows) >= 1:
                last_idx = len(target.rows) - 1
                last_row = target.rows[last_idx]
                # R[-1]C0 = "Total", R[-1]C3 = "100%"
                for c_idx, val in [(0, 'Total'), (3, '100%')]:
                    if c_idx < len(target.columns):
                        cell = last_row.cells[c_idx]
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                run.text = ''
                        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                            cell.text_frame.paragraphs[0].runs[0].text = val

    # 2x1 표 — 헤더에 따라 각각 다른 내용 채우기
    for tbl_shape in tables:
        t = tbl_shape.table
        if len(t.rows) == 2 and len(t.columns) == 1:
            header = t.rows[0].cells[0].text_frame.text.strip()
            cell = t.rows[1].cells[0]

            if '경쟁사' in header:
                content = (
                    '럭셔리 뷰티 11개 브랜드: Dior, Chanel, Hera, MAC, Jo Malone, '
                    'Nars, Estee Lauder, Lancome, Tom Ford, Sulwhasoo, Prada Beauty'
                )
            elif '분석' in header and '설계' in header:
                content = (
                    '3 카테고리(향수/기프팅/쿠션) × 4 인텐트(니즈 인식/정보 탐색/대안 비교/구매 결정) 매트릭스. '
                    '카테고리별 40/40/20 비율로 360개 제네릭 질문 생성, 키워드 360개와 1:1 매핑. '
                    '3개 채널(ChatGPT/Google AIO/Naver) × 4 사이클 누적 응답을 자사·경쟁사 11개 텍스트 매칭으로 분석.'
                )
            else:
                continue

            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ''
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                cell.text_frame.paragraphs[0].runs[0].text = content
            elif cell.text_frame.paragraphs:
                cell.text_frame.paragraphs[0].add_run().text = content


def _fill_slide_8_funnel(slide):
    """슬라이드 8: 퍼널 4표 - 전체 / 채널 / 카테고리 / 인텐트."""
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) < 4:
        return

    # 1. 전체 퍼널 (table 0): 7x4 — Total → AI Existence → Brand Mention → Citation
    f = metrics['funnel_all']
    total_cit = metrics['total_citations']
    ysl_cit = metrics['ysl_brand_url_count']
    overall_data = [
        ['Total Query Set', f"{f['total']:,}", '100%', '쿼리 (질문 + 키워드) × 3 AI Engine × 4 사이클'],
        ['AI Existence', f"{f['ai_existence']:,}", f"{f['ai_existence_rate']:.1f}%", '응답 생성된 쿼리'],
        ['Brand Mention', f"{f['brand_mention']:,}", f"{f['brand_mention_rate']:.1f}%", 'YSL 언급 응답'],
        ['Citation (전체)', f"{total_cit:,}", '—', '응답 내 인용된 URL 총 수'],
        ['Citation (YSL 자사)', f"{ysl_cit}", f"{ysl_cit/total_cit*100:.2f}%", 'YSL 자사몰 도메인 인용'],
    ]
    fill_table(tables[0].table, overall_data, start_row=1, clear_extra_cols=True)

    # 2. 채널별 (table 1): 7x4 — Channel / 응답 / 멘션 / 멘션률
    chan_data = []
    for r in metrics['by_channel']:
        chan_data.append([
            CHANNEL_LABEL.get(r['channel'], r['channel']),
            f"{r['total']:,}",
            f"{r['mention']:,}",
            f"{r['rate']:.2f}%",
        ])
    fill_table(tables[1].table, chan_data, start_row=1, clear_extra_cols=True)
    # R0 (타이틀, 병합셀) 복원
    cell = tables[1].table.rows[0].cells[0]
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.text = ''
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        cell.text_frame.paragraphs[0].runs[0].text = 'B. Channel Funnel – AI 플랫폼 별 Visibility'
    elif cell.text_frame.paragraphs:
        cell.text_frame.paragraphs[0].add_run().text = 'B. Channel Funnel – AI 플랫폼 별 Visibility'

    # 3. 카테고리별 (table 2): 7x7 — 카테고리 / 응답 / 멘션 / 멘션률 / 인용수 / 인용 점유율 / 응답당 인용
    cat_data = []
    total_cit_all = metrics['total_citations']
    for r in metrics['by_category']:
        cat = r['category']
        cit = metrics['cit_by_category'].get(cat, 0)
        cit_share = cit / total_cit_all * 100 if total_cit_all else 0
        cit_per_resp = cit / r['total'] if r['total'] else 0
        cat_data.append([
            cat,
            f"{r['total']:,}",
            f"{r['mention']:,}",
            f"{r['rate']:.2f}%",
            f"{cit:,}",
            f"{cit_share:.1f}%",
            f"{cit_per_resp:.1f}",
        ])
    fill_table(tables[2].table, cat_data, start_row=1, clear_extra_cols=True)
    # 헤더 R0: 7컬럼 채움
    cat_headers = ['Category', '응답수', '멘션수', '멘션률', '인용수', '인용률', '비고']
    cell = tables[2].table.rows[0].cells[0]
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.text = ''
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        cell.text_frame.paragraphs[0].runs[0].text = 'C. Category Funnel – 향수/기프팅/쿠션 카테고리별 Visibility'
    elif cell.text_frame.paragraphs:
        cell.text_frame.paragraphs[0].add_run().text = 'C. Category Funnel – 향수/기프팅/쿠션 카테고리별 Visibility'
    # 단, 7컬럼이지만 R0C0이 타이틀이라 cols 1~ 헤더는 R1부터? — 아니, 원본 SKC에서 R0C0은 merged cell. R1이 헤더, R2부터 데이터.
    # 다시 살펴 — 현재 fill은 R1부터 데이터 들어가는 구조. 헤더는 cols 1~6에 R0이 합쳐져 있음.
    # 안전하게: R1이 헤더가 아니라면 그냥 R0~ 데이터 진행. 추가 컬럼 헤더는 cell이 있으면 채움
    for c, h in enumerate(cat_headers):
        if c == 0: continue  # R0C0은 타이틀
        if c >= len(tables[2].table.columns):
            break

    # 4. 인텐트별 (table 3): 7x6 — 인텐트 / 응답 / 멘션 / 멘션률 / 인용수 / 응답당 인용
    int_data = []
    for r in metrics['by_intent']:
        intent = r['intent']
        cit = metrics['cit_by_intent'].get(intent, 0)
        cit_per_resp = cit / r['total'] if r['total'] else 0
        int_data.append([
            intent,
            f"{r['total']:,}",
            f"{r['mention']:,}",
            f"{r['rate']:.2f}%",
            f"{cit:,}",
            f"{cit_per_resp:.1f}",
        ])
    fill_table(tables[3].table, int_data, start_row=1, clear_extra_cols=True)


def _fill_slide_9_citation(slide):
    """슬라이드 9: 인용 표 — Top 도메인 (6x4) + YSL e-commerce 채널 (9x5)."""
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) < 2:
        return

    # 표 #0 (6x4): No / 도메인 / 인용수 / 점유율 — 헤더부터 다시 채움
    t0 = tables[0].table
    _set_cell(t0.rows[0].cells[0], 'No.')
    _set_cell(t0.rows[0].cells[1], '도메인')
    _set_cell(t0.rows[0].cells[2], '인용 수')
    _set_cell(t0.rows[0].cells[3], '점유율')

    top = list(metrics['top_domains_all'].items())[:5]
    total_cit = metrics['total_citations']
    domain_data = [
        [str(i + 1), d, f"{c:,}", f"{c/total_cit*100:.1f}%"]
        for i, (d, c) in enumerate(top)
    ]
    fill_table(t0, domain_data, start_row=1, clear_extra_cols=True)

    # 표 #1 (9x5): YSL E-Commerce Channel mapping
    # 헤더: No. / 채널 / Target URL / 페이지 설명 / 인용 결과
    t1 = tables[1].table
    _set_cell(t1.rows[0].cells[0], 'No.')
    _set_cell(t1.rows[0].cells[1], '채널')
    _set_cell(t1.rows[0].cells[2], 'Target URL')
    _set_cell(t1.rows[0].cells[3], '페이지 설명')
    _set_cell(t1.rows[0].cells[4], '인용 결과')

    # 인용 도메인을 lower-case로 검사
    cit_domains_lower = citation_df['Domain'].astype(str).str.lower()
    def cited(*patterns):
        for p in patterns:
            n = cit_domains_lower.str.contains(p, regex=False, na=False).sum()
            if n > 0:
                return f'인용 {int(n)}건'
        return '미인용'

    ec_channels = [
        ('YSL 자사몰', 'yslbeauty.co.kr', '브랜드 공식몰', cited('ysl', 'yslbeauty')),
        ('네이버 브랜드 스토어', 'brand.naver.com/yslbeauty', '네이버 공식 브랜드 스토어', cited('brand.naver.com', 'smartstore.naver.com/ysl')),
        ('롯데온', 'lotteon.com', '롯데 종합몰', cited('lotteon')),
        ('쿠팡 YSL', 'coupang.com (YSL 검색)', '쿠팡 셀러몰', cited('coupang.com')),
        ('카카오 선물하기', 'gift.kakao.com', '모바일 기프팅', cited('gift.kakao.com')),
        ('SSG/신세계', 'ssg.com', '백화점 온라인', cited('ssg.com', 'shinsegae')),
        ('11번가', '11st.co.kr', '오픈마켓', cited('11st')),
        ('올리브영', 'oliveyoung.co.kr', '뷰티 전문몰 (YSL 미입점)', cited('oliveyoung')),
    ]
    rows_data = [
        [str(i + 1), name, url, desc, result]
        for i, (name, url, desc, result) in enumerate(ec_channels)
    ]
    fill_table(t1, rows_data, start_row=1, clear_extra_cols=True)

    # 텍스트박스 "Skinceuticals domain citation" → "YSL domain citation"
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.has_table:
            if 'Skinceuticals' in shape.text_frame.text or 'skinceuticals' in shape.text_frame.text.lower():
                _set_cell_textframe = shape.text_frame
                for para in _set_cell_textframe.paragraphs:
                    for run in para.runs:
                        run.text = ''
                if _set_cell_textframe.paragraphs and _set_cell_textframe.paragraphs[0].runs:
                    _set_cell_textframe.paragraphs[0].runs[0].text = 'YSL Top Citation Domains'
                elif _set_cell_textframe.paragraphs:
                    _set_cell_textframe.paragraphs[0].add_run().text = 'YSL Top Citation Domains'


def _unmerge_cell(cell):
    """vMerge / gridSpan 속성 제거하여 셀 병합 해제."""
    tc = cell._tc
    if tc.get('vMerge') is not None:
        del tc.attrib['vMerge']
    if tc.get('gridSpan') is not None:
        del tc.attrib['gridSpan']
    if tc.get('rowSpan') is not None:
        del tc.attrib['rowSpan']
    # hMerge도 제거
    if tc.get('hMerge') is not None:
        del tc.attrib['hMerge']


def _fill_slide_11_competitor(slide):
    """슬라이드 11: 채널 × 경쟁사 Mention Rate (6x19)."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return
    table = tables[0].table

    # 모든 셀의 병합 속성 제거 (R1 vMerge로 ChatGPT 행이 가려지는 문제 해결)
    for row in table.rows:
        for cell in row.cells:
            _unmerge_cell(cell)

    # 컬럼 폭 재배분 — 14 컬럼만 사용, 나머지 5개는 폭 0으로 숨김
    # 합계 12.02" 유지
    widths_in = [
        1.2,   # 0: Channel
        0.95,  # 1: Total Responses
        0.95,  # 2: YSL Beauty
        0.74, 0.74, 0.74, 0.74, 0.74, 0.74, 0.74, 0.74, 0.74, 0.74, 0.74,  # 3~13: 11 경쟁사 (Dior~Prada)
        0, 0, 0, 0, 0,  # 14~18: 미사용
    ]
    for c_idx, w in enumerate(widths_in):
        if c_idx < len(table.columns):
            table.columns[c_idx].width = Emu(int(w * 914400))

    brands_order = ['YSL Beauty', 'Dior', 'Chanel', 'Hera', 'Jo Malone', 'Estee Lauder',
                    'Lancome', 'Tom Ford', 'Sulwhasoo', 'MAC', 'Nars', 'Prada Beauty']

    # 헤더 R0 채우기 (Channel / Total / 12 브랜드)
    headers = ['Channel', 'Total Responses'] + brands_order
    for c, h in enumerate(headers):
        if c >= len(table.columns):
            break
        cell = table.rows[0].cells[c]
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''
        if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
            cell.text_frame.paragraphs[0].runs[0].text = h
        elif cell.text_frame.paragraphs:
            cell.text_frame.paragraphs[0].add_run().text = h
    # R0 추가 컬럼 클리어
    for c in range(len(headers), len(table.columns)):
        cell = table.rows[0].cells[c]
        for para in cell.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''

    rows_data = []
    for chan_row in metrics['sov_by_channel']:
        chan = chan_row.get('channel', '')
        row = [CHANNEL_LABEL.get(chan, chan)]
        row.append(f"{chan_row.get('total', 0):,}")
        for b in brands_order:
            v = chan_row.get(b, 0)
            row.append(f"{v:.1f}%")
        rows_data.append(row)
    fill_table(table, rows_data, start_row=1, clear_extra_cols=True)


def _set_cell(cell, value):
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.text = ''
    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
        cell.text_frame.paragraphs[0].runs[0].text = value
    elif cell.text_frame.paragraphs:
        cell.text_frame.paragraphs[0].add_run().text = value


def _fill_slide_12_engine_compare(slide):
    """슬라이드 12: 엔진별 비교 (4x4) + 카테고리 spotlight (5x4) + 브랜드 비교 (2x9)."""
    tables = [s for s in slide.shapes if s.has_table]

    sov_overall = metrics['sov_overall']
    sov_by_channel = metrics['sov_by_channel']
    sov_by_category = metrics['sov_by_category']

    for tbl_shape in tables:
        t = tbl_shape.table
        rows, cols = len(t.rows), len(t.columns)

        # 4x4 표: Mention rate by engine — Engine / YSL / 경쟁평균 / Gap
        if rows == 4 and cols == 4:
            # Header
            _set_cell(t.rows[0].cells[0], 'Engine')
            _set_cell(t.rows[0].cells[1], 'YSL Beauty (%)')
            _set_cell(t.rows[0].cells[2], '경쟁사 평균 (%)')
            _set_cell(t.rows[0].cells[3], 'Gap (%p)')

            comp_brands = [b for b in sov_overall.keys() if b != 'YSL Beauty']
            engine_rows = [(c, CHANNEL_LABEL.get(c, c)) for c in ['chatgpt', 'google', 'naver']]
            for r_idx, (chan, label) in enumerate(engine_rows, start=1):
                if r_idx >= rows: break
                chan_row = next((cr for cr in sov_by_channel if cr['channel'] == chan), None)
                if not chan_row: continue
                ysl_v = chan_row.get('YSL Beauty', 0)
                comp_avg = sum(chan_row.get(b, 0) for b in comp_brands) / len(comp_brands)
                gap = ysl_v - comp_avg
                _set_cell(t.rows[r_idx].cells[0], label)
                _set_cell(t.rows[r_idx].cells[1], f'{ysl_v:.1f}%')
                _set_cell(t.rows[r_idx].cells[2], f'{comp_avg:.1f}%')
                _set_cell(t.rows[r_idx].cells[3], f'{gap:+.1f}%p')

        # 5x4 표: Category spotlight — Category / YSL / 경쟁평균 / Gap
        elif rows == 5 and cols == 4:
            _set_cell(t.rows[0].cells[0], 'Category')
            _set_cell(t.rows[0].cells[1], 'YSL Beauty (%)')
            _set_cell(t.rows[0].cells[2], '경쟁사 평균 (%)')
            _set_cell(t.rows[0].cells[3], 'Gap (%p)')

            comp_brands = [b for b in sov_overall.keys() if b != 'YSL Beauty']
            for r_idx, cat in enumerate(['향수', '기프팅', '쿠션'], start=1):
                if r_idx >= rows: break
                cat_row = next((cr for cr in sov_by_category if cr['category'] == cat), None)
                if not cat_row: continue
                ysl_v = cat_row.get('YSL Beauty', 0)
                comp_avg = sum(cat_row.get(b, 0) for b in comp_brands) / len(comp_brands)
                gap = ysl_v - comp_avg
                _set_cell(t.rows[r_idx].cells[0], cat)
                _set_cell(t.rows[r_idx].cells[1], f'{ysl_v:.1f}%')
                _set_cell(t.rows[r_idx].cells[2], f'{comp_avg:.1f}%')
                _set_cell(t.rows[r_idx].cells[3], f'{gap:+.1f}%p')
            # R4 클리어
            if rows >= 5:
                for c in range(cols):
                    for para in t.rows[4].cells[c].text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''

        # 2x9 표: Mention rate by brand — YSL Beauty + 경쟁사 7개
        elif rows == 2 and cols >= 8:
            brands = ['YSL Beauty', 'Dior', 'Chanel', 'Hera', 'Jo Malone',
                      'Estee Lauder', 'Lancome', 'Sulwhasoo']
            _set_cell(t.rows[0].cells[0], 'Metric')
            for c, b in enumerate(brands, start=1):
                if c >= cols: break
                _set_cell(t.rows[0].cells[c], b)
            _set_cell(t.rows[1].cells[0], 'Overall Mention Rate')
            for c, b in enumerate(brands, start=1):
                if c >= cols: break
                _set_cell(t.rows[1].cells[c], f"{sov_overall.get(b, 0):.1f}%")
            # 추가 컬럼 클리어
            for c in range(len(brands) + 1, cols):
                for r_idx in range(rows):
                    for para in t.rows[r_idx].cells[c].text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''


def _fill_slide_13_intent_journey(slide):
    """슬라이드 13: 인텐트 저니별 Mention Rate (14x8 표)."""
    tables = [s for s in slide.shapes if s.has_table]
    # 14x8 표 찾기
    for tbl_shape in tables:
        t = tbl_shape.table
        if len(t.rows) >= 10 and len(t.columns) >= 6:
            # 인텐트 4개 × 카테고리 3개 매트릭스
            intents = ['니즈 인식', '정보 탐색', '대안 비교', '구매 결정']
            cats = ['향수', '기프팅', '쿠션', '전체']

            # 헤더 (R0): 인텐트 / 향수 / 기프팅 / 쿠션 / 전체 / [TBD]
            headers = ['인텐트 (Customer Journey)', '향수', '기프팅', '쿠션', 'Total', 'YSL Mention']
            for c, h in enumerate(headers):
                if c >= len(t.columns):
                    break
                cell = t.rows[0].cells[c]
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''
                if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].text = h
                elif cell.text_frame.paragraphs:
                    cell.text_frame.paragraphs[0].add_run().text = h

            # 인텐트별 행 (R1~R4)
            for r_idx, intent in enumerate(intents, start=1):
                if r_idx >= len(t.rows):
                    break
                # C0: intent name
                cell = t.rows[r_idx].cells[0]
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''
                if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                    cell.text_frame.paragraphs[0].runs[0].text = intent
                # C1~C3: 카테고리별 YSL mention rate
                for c_idx, cat in enumerate(['향수', '기프팅', '쿠션'], start=1):
                    sub = mention_df[(mention_df['intent'] == intent) & (mention_df['category'] == cat)]
                    n = len(sub)
                    rate = sub['YSL Beauty'].sum() / max(n, 1) * 100 if n > 0 else 0
                    cell = t.rows[r_idx].cells[c_idx]
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''
                    val = f'{rate:.1f}%' if n >= 10 else f'-'
                    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].text = val
                # C4: Total
                sub = mention_df[mention_df['intent'] == intent]
                rate_total = sub['YSL Beauty'].sum() / max(len(sub), 1) * 100 if len(sub) > 0 else 0
                if 4 < len(t.columns):
                    cell = t.rows[r_idx].cells[4]
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''
                    if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                        cell.text_frame.paragraphs[0].runs[0].text = f'{rate_total:.1f}%'
                # 나머지 컬럼 클리어
                for c in range(5, len(t.columns)):
                    cell = t.rows[r_idx].cells[c]
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''

            # 5행 이후 클리어 (SKC 추가 row)
            for r_idx in range(5, len(t.rows)):
                for c in range(len(t.columns)):
                    cell = t.rows[r_idx].cells[c]
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = ''
            break


def _fill_slide_14_category_deep_dive(slide):
    """슬라이드 14: 카테고리별 브랜드 랭킹 표 3개 (향수/기프팅/쿠션)."""
    tables = [s for s in slide.shapes if s.has_table]

    cats = ['향수', '기프팅', '쿠션']
    cat_titles = {
        '향수': '향수 카테고리: 브랜드 Mention 순위',
        '기프팅': '기프팅 카테고리: 브랜드 Mention 순위',
        '쿠션': '쿠션 카테고리: 브랜드 Mention 순위',
    }
    # 카테고리별 핵심 인사이트 (YSL 기준)
    cat_notes = {
        '향수': {
            'Dior': '여성/남성 향수 전반 1위 노출',
            'Jo Malone': '럭셔리 + 입문 향수에서 강세',
            'Chanel': '시그니처 라인 (No.5/Coco) 인지',
            'Tom Ford': '니치/프리미엄 영역',
            'YSL Beauty': 'Libre/Black Opium 일부 노출 — 브랜드 다양성 약함',
        },
        '기프팅': {
            'Dior': '여친/엄마 선물 전반 1위',
            'Sulwhasoo': '40~60대 어머니 선물 강세',
            'Chanel': '20~30대 여성 선물',
            'YSL Beauty': '카카오 선물하기 채널 강세',
            'Lancome': '면세점 선물 영역',
        },
        '쿠션': {
            'Hera': '쿠션 카테고리 명품 1위 (블랙쿠션)',
            'YSL Beauty': '명품 쿠션 2위 — Touche Éclat Le Cushion',
            'Dior': '럭셔리 쿠션 (Forever Cushion)',
            'Chanel': 'Le Teint Ultra Cushion',
            'MAC': '커버력 쿠션',
        },
    }

    for i, cat in enumerate(cats[:len(tables)]):
        cat_row = next((r for r in metrics['sov_by_category'] if r.get('category') == cat), None)
        if not cat_row:
            continue
        tbl = tables[i].table
        # R0: 카테고리 타이틀 (병합 셀)
        _set_cell(tbl.rows[0].cells[0], cat_titles[cat])

        # R1: 헤더 — 4컬럼만 사용 + 5/6컬럼 클리어
        if len(tbl.rows) >= 2 and len(tbl.columns) >= 4:
            headers = ['순위', '브랜드', 'Mention Rate', '비고']
            for c, h in enumerate(headers):
                if c >= len(tbl.columns):
                    break
                _set_cell(tbl.rows[1].cells[c], h)
            for c in range(len(headers), len(tbl.columns)):
                cell = tbl.rows[1].cells[c]
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''

        # 브랜드 랭킹 — 표 행 수에 맞춰 채움 (모든 빈 행을 데이터로 채워서 stretched 해소)
        brand_rates = [(b, cat_row.get(b, 0)) for b in metrics['sov_overall'].keys() if b in cat_row]
        brand_rates.sort(key=lambda x: -x[1])
        max_data_rows = len(tbl.rows) - 2  # R0=title, R1=header, R2~=data
        top_data = []
        for j, (b, r) in enumerate(brand_rates[:max_data_rows]):
            note = cat_notes.get(cat, {}).get(b, '')
            top_data.append([j + 1, b, f"{r:.2f}%", note])
        fill_table(tbl, top_data, start_row=2, clear_extra_cols=True)


def _fill_slide_17_positioning(slide):
    """슬라이드 17: 표 (5x7)."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return
    # 인지 = 정보탐색+니즈인식 / 전환 = 대안비교+구매결정
    upper_intents = ['니즈 인식', '정보 탐색']
    lower_intents = ['대안 비교', '구매 결정']

    upper = mention_df[mention_df['intent'].isin(upper_intents)]
    lower = mention_df[mention_df['intent'].isin(lower_intents)]

    rows_data = []
    for b in ['YSL Beauty', 'Dior', 'Chanel', 'Hera', 'Jo Malone', 'Estee Lauder']:
        if b in mention_df.columns:
            u = upper[b].sum() / max(len(upper), 1) * 100
            l = lower[b].sum() / max(len(lower), 1) * 100
            rows_data.append([b, f"{u:.2f}%", f"{l:.2f}%", '', ''])

    fill_table(tables[0].table, rows_data, start_row=1, clear_extra_cols=True)


def _fill_slide_17_bubbles(slide):
    """슬라이드 17 포지셔닝 맵의 버블 라벨을 YSL 경쟁사로 교체."""
    brand_counts = {}
    for b in mention_df.columns:
        if b in metrics['sov_overall'].keys():
            brand_counts[b] = int(mention_df[b].sum())

    brand_list = sorted(brand_counts.items(), key=lambda x: -x[1])

    # 1. SKC 잔존 브랜드 라벨 (이니스프리/이지듀/설화수/토리든 등) 클리어
    SKC_BRANDS_KR = ['이니스프리', '이지듀', '설화수', '토리든', '아이오페', '클레어스', 'SK-II', 'SK2', '랑콤', '에스티로더']
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.has_table:
            t = shape.text_frame.text.strip()
            if any(kw in t for kw in SKC_BRANDS_KR) and re.search(r'\(\d+\)', t):
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        run.text = ''

    # 2. "(숫자)" 패턴 또는 빈 텍스트박스를 YSL 경쟁사 라벨로 순차 채움
    label_idx = 0
    for shape in slide.shapes:
        if shape.has_text_frame and not shape.has_table:
            t = shape.text_frame.text.strip()
            if re.match(r'^[가-힣A-Za-z]+\s*\(\d+\)$', t) or t == '':
                if label_idx < len(brand_list):
                    b, c = brand_list[label_idx]
                    label_idx += 1
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = ''
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        tf.paragraphs[0].runs[0].text = f'{b} ({c})'
                    elif tf.paragraphs:
                        run = tf.paragraphs[0].add_run()
                        run.text = f'{b} ({c})'


def _fill_slide_19_sov(slide):
    """슬라이드 19: SOV 랭킹 3표 (전체 / 향수 / 쿠션 등)."""
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) < 3:
        return

    # 표 1: 전체 SOV
    sov_sorted = sorted(metrics['sov_overall'].items(), key=lambda x: -x[1])[:10]
    sov_data = [[i + 1, b, f"{r:.2f}%", '', ''] for i, (b, r) in enumerate(sov_sorted)]
    fill_table(tables[0].table, sov_data, start_row=1)

    # 표 2: 카테고리별 (향수)
    perfume_row = next((r for r in metrics['sov_by_category'] if r.get('category') == '향수'), None)
    if perfume_row:
        brands_sorted = sorted(
            [(b, perfume_row.get(b, 0)) for b in metrics['sov_overall'].keys()],
            key=lambda x: -x[1],
        )[:10]
        perfume_data = [[i + 1, b, f"{r:.2f}%", '', ''] for i, (b, r) in enumerate(brands_sorted)]
        fill_table(tables[1].table, perfume_data, start_row=1)

    # 표 3: 카테고리별 (기프팅)
    gift_row = next((r for r in metrics['sov_by_category'] if r.get('category') == '기프팅'), None)
    if gift_row:
        brands_sorted = sorted(
            [(b, gift_row.get(b, 0)) for b in metrics['sov_overall'].keys()],
            key=lambda x: -x[1],
        )[:10]
        gift_data = [[i + 1, b, f"{r:.2f}%", '', ''] for i, (b, r) in enumerate(brands_sorted)]
        fill_table(tables[2].table, gift_data, start_row=1)


def _fill_slide_23_citation_charts(slide):
    """슬라이드 23: 도메인 / 콘텐츠 타입 차트 2개."""
    charts = [s for s in slide.shapes if s.has_chart]
    if not charts:
        return

    # 차트 1: 도메인 타입 분포
    if len(charts) >= 1:
        cd1 = CategoryChartData()
        types = sorted(metrics['domain_type_dist'].items(), key=lambda x: -x[1])[:8]
        cd1.categories = [t for t, _ in types]
        cd1.add_series('인용 비중', [c for _, c in types])
        try:
            charts[0].chart.replace_data(cd1)
        except Exception as e:
            print(f'슬라이드 23 차트1 교체 실패: {e}')

    # 차트 2: Top 도메인
    if len(charts) >= 2:
        cd2 = CategoryChartData()
        top = list(metrics['top_domains_all'].items())[:8]
        cd2.categories = [d for d, _ in top]
        cd2.add_series('인용수', [c for _, c in top])
        try:
            charts[1].chart.replace_data(cd2)
        except Exception as e:
            print(f'슬라이드 23 차트2 교체 실패: {e}')


def _fill_slide_24_top_domains(slide):
    """슬라이드 24: 채널별 Top 도메인 표 3개 (ChatGPT / Google / Naver)."""
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) < 3:
        return

    plats = ['chatgpt', 'google', 'naver']
    for i, plat in enumerate(plats[:len(tables)]):
        if plat in metrics['top_domains_by_platform']:
            domains = list(metrics['top_domains_by_platform'][plat].items())[:7]
            total = sum(metrics['top_domains_by_platform'][plat].values())
            data = [[j + 1, d, f"{c}", f"{c/max(total,1)*100:.1f}%"] for j, (d, c) in enumerate(domains)]
            fill_table(tables[i].table, data, start_row=1)


def _fill_slide_26_ecommerce(slide):
    """슬라이드 26: 이커머스 도메인 인용."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return

    ecom = citation_df[citation_df['domain_type'] == 'Commerce']
    top = ecom['Domain'].value_counts().head(10)
    total_ecom = len(ecom)
    data = [[i + 1, d, f"{c}", f"{c/max(total_ecom,1)*100:.1f}%", ''] for i, (d, c) in enumerate(top.items())]
    fill_table(tables[0].table, data, start_row=1)

    # 표 #1: 이커머스 인용 특징 요약 — 헤더 두고 데이터 행은 클리어 + TBD
    if len(tables) >= 2:
        t = tables[1].table
        for r_idx in range(2, len(t.rows)):
            for c_idx in range(len(t.columns)):
                cell = t.rows[r_idx].cells[c_idx]
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = ''
        # R2C0에 TBD 표시
        if len(t.rows) >= 3:
            cell = t.rows[2].cells[0]
            if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs:
                cell.text_frame.paragraphs[0].runs[0].text = '[TBD] 페이지 유형 라벨링 필요'
            elif cell.text_frame.paragraphs:
                cell.text_frame.paragraphs[0].add_run().text = '[TBD] 페이지 유형 라벨링 필요'


def _fill_slide_15_persona(slide):
    """슬라이드 15: 페르소나·스킨타입 Mention Rate (축2 기반)."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return

    # 축2 라벨별 응답 매핑 (Reference ID → axis2)
    axis2_map = dict(zip(qmeta_df['no'].astype(str), qmeta_df['axis2']))

    def parse_no(ref):
        if not isinstance(ref, str):
            return None
        m = re.search(r'-(\d+)$', ref)
        return m.group(1).lstrip('0') if m else None

    mention_df['no_extracted'] = mention_df['reference_id'].astype(str).apply(parse_no)
    mention_df['axis2'] = mention_df['no_extracted'].map(axis2_map)

    persona_groups = ['① 연령/취향/입문', '② 연령대별', '③ 피부타입', '④ 남성 구매자', '⑤ 남성 본인']

    # 표 1 (6x9): 페르소나 × 표본수·YSL Rate·Top 경쟁사·신뢰도
    rows_data1 = []
    for p in persona_groups:
        sub = mention_df[mention_df['axis2'] == p]
        n = len(sub)
        ysl_rate = sub['YSL Beauty'].sum() / max(n, 1) * 100 if n > 0 else 0
        comp_rates = []
        for b in ['Dior', 'Chanel', 'Jo Malone', 'Hera', 'Sulwhasoo']:
            if b in sub.columns and n > 0:
                comp_rates.append((b, sub[b].sum() / n * 100))
        comp_rates.sort(key=lambda x: -x[1])
        top1 = comp_rates[0] if comp_rates else ('-', 0)
        rows_data1.append([
            p, f'n={n}', f'{ysl_rate:.1f}%',
            f'{top1[0]} {top1[1]:.1f}%',
            '⚠️ 표본 부족' if n < 30 else 'OK',
        ])
    fill_table(tables[0].table, rows_data1, start_row=1)

    # 표 2 (8x9): 페르소나별 × YSL + 주요 경쟁사 5개 Mention Rate
    if len(tables) >= 2:
        rows_data2 = []
        comp_brands = ['Dior', 'Chanel', 'Jo Malone', 'Hera', 'Sulwhasoo']
        for p in persona_groups:
            sub = mention_df[mention_df['axis2'] == p]
            n = len(sub)
            row = [p, f'{n}', f"{sub['YSL Beauty'].sum()/max(n,1)*100:.1f}%"]
            for b in comp_brands:
                if b in sub.columns and n > 0:
                    row.append(f"{sub[b].sum()/n*100:.1f}%")
                else:
                    row.append('0.0%')
            rows_data2.append(row)
        fill_table(tables[1].table, rows_data2, start_row=1)


def _fill_slide_27_brand_mall(slide):
    """슬라이드 27: 자사몰 인용."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return

    # YSL 자사몰
    ysl_url = citation_df[citation_df['Domain'].str.contains('ysl', case=False, na=False)]
    rows_data = [
        [1, 'yslbeauty.co.kr', f"{len(ysl_url)}", '', '', '', ''],
        [2, 'YSL 외 브랜드 자사몰', '0', '', '', '', ''],
    ]
    fill_table(tables[0].table, rows_data, start_row=1, clear_extra_cols=True)


def _fill_slide_32_analysis_angle(slide):
    """슬라이드 32: 분석 각도 (축1/2/3 기반 8x4 표)."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return

    # 축별 분포 카운트
    a1 = qmeta_df['axis1'].value_counts()
    a2 = qmeta_df['axis2'].value_counts()
    a3 = qmeta_df['axis3'].value_counts()

    rows_data = [
        ['축1: 사용 상황/TPO', '계절·기념일·선물 상황별 탐색 패턴', '계절/TPO, 기념일/이벤트, 선물 상황', f'{(360 - a1.get("-", 0))}'],
        ['축2: 소비자 프로필', '연령·피부타입·구매자 관점별 차이', '연령/취향/입문, 피부타입, 남성 구매자/본인', f'{(360 - a2.get("-", 0))}'],
        ['축3: 구매 허들', '실제 구매 의사결정 시 장애 요소', '품질, 브랜드, 가격, 채널, 실패 방지 등 8개', f'{(360 - a3.get("-", 0))}'],
        ['카테고리', '뷰티 핵심 카테고리 3개', '향수(40%), 기프팅(40%), 쿠션(20%)', '360'],
        ['Intent', '소비자 의사결정 4단계', '니즈 인식 / 정보 탐색 / 대안 비교 / 구매 결정', '360'],
        ['Question', '제네릭 자연어 질문', '브랜드명 배제, 존댓말, ?/. 종결', '360'],
        ['Keyword', '1:1 검색 키워드 매핑', '질문별 unique 검색어 (공백 없는 한국어)', '360'],
    ]
    fill_table(tables[0].table, rows_data, start_row=1, clear_extra_cols=True)


def _fill_slide_33_intent_def(slide):
    """슬라이드 33: 인텐트 정의 표."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return

    rows_data = []
    for d in INTENT_DEFS:
        if not d.get('name'):
            continue
        rows_data.append([
            (d['name'] or '').replace('\n', ' '),
            (d['def'] or '').replace('\n', ' ')[:100],
            (d['examples'] or '').replace('\n', ' / ')[:100],
            f"{(d['ratio'] or 0)*100:.0f}%",
        ])
    # 남는 행 클리어 (SKC의 How-to 등 잔존 방지)
    while len(rows_data) < 6:
        rows_data.append(['', '', '', ''])
    fill_table(tables[0].table, rows_data, start_row=1, clear_extra_cols=True)


def _fill_slide_34_benchmark(slide):
    """슬라이드 34: KW vs Q 벤치마크."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return

    # KW 응답 vs Q 응답 멘션률 비교
    kw_only = mention_df[mention_df['channel'].isin(['google', 'naver'])]
    q_only = mention_df[mention_df['channel'] == 'chatgpt']

    def m_rate(df, brand):
        n = len(df)
        return df[brand].sum() / max(n, 1) * 100 if brand in df.columns and n > 0 else 0

    rows_data = []
    for b in ['YSL Beauty', 'Dior', 'Chanel', 'Hera', 'Jo Malone', 'Lancome']:
        kw_r = m_rate(kw_only, b)
        q_r = m_rate(q_only, b)
        gap = abs(kw_r - q_r)
        rows_data.append([
            b,
            f'{kw_r:.1f}%',
            f'{q_r:.1f}%',
            f'{gap:.1f}%p',
        ])
    fill_table(tables[0].table, rows_data, start_row=1, clear_extra_cols=True)


# =============================================================================
# 메인
# =============================================================================

def delete_slides(prs, indices_1based):
    """python-pptx로 슬라이드 삭제 (XML 직접 조작)."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    # 뒤에서부터 삭제 (인덱스 안 밀림)
    for idx in sorted(indices_1based, reverse=True):
        sld = slides[idx - 1]
        rId = sld.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)


def update_toc(slide):
    """슬라이드 2: 목차 — 4섹션 구조로 갱신."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return
    t = tables[0].table  # 6x3
    # SKC 원본은 6행 (5섹션 + appendix). YSL은 4섹션만 사용.
    toc_data = [
        ['1', 'AI 가시성 현황 진단', '쿼리 4,320 / Mention 9.6% / 자사몰 인용 0.11% — AI 답변에서 YSL 노출 수준 정량화'],
        ['2', '포지셔닝 검증', '엔진별·카테고리별 경쟁사 비교 — Dior 22%·Chanel 17%·YSL 9.6%, 향수에서 -2%p 열세'],
        ['3', 'Share of Voice', '카테고리별 SOV 랭킹 — 쿠션 +7.3%p 우위, 향수 보강 우선'],
        ['4', '콘텐츠 인용 소스 구조', 'Top 도메인 / 채널별 인용 패턴 / 자사몰 12건 (0.11%)'],
    ]
    # 데이터 채움 (R1~R4)
    for i, row_data in enumerate(toc_data):
        r = i + 1
        if r >= len(t.rows): break
        for c, val in enumerate(row_data):
            if c >= len(t.columns): break
            _set_cell(t.rows[r].cells[c], val)
    # 5~ 행 클리어 (SKC의 section 5 + Appendix)
    for r in range(len(toc_data) + 1, len(t.rows)):
        for c in range(len(t.columns)):
            cell = t.rows[r].cells[c]
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ''


def unmerge_all_tables(prs):
    """모든 슬라이드의 모든 표 셀에서 vMerge/gridSpan/rowSpan 속성 제거.
    vMerge로 데이터 행이 가려지는 문제 방지.
    """
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        tc = cell._tc
                        if tc.get('vMerge') is not None:
                            del tc.attrib['vMerge']
                            count += 1
                        if tc.get('gridSpan') is not None:
                            del tc.attrib['gridSpan']
                            count += 1
                        if tc.get('rowSpan') is not None:
                            del tc.attrib['rowSpan']
                        if tc.get('hMerge') is not None:
                            del tc.attrib['hMerge']
    print(f'병합 속성 제거: {count}건')


def replace_skc_logo(prs):
    """모든 슬라이드의 SKC 로고를 YSL 로고로 교체.
    SKC 로고는 슬라이드 우측 상단 (left>9", top<1") 또는 1번 슬라이드 하단의 작은 picture."""
    YSL_LOGO = ROOT / 'assets' / 'ysl_logo_1.png'
    if not YSL_LOGO.exists():
        print(f'YSL 로고 파일 없음: {YSL_LOGO}')
        return

    replaced = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        # 모든 picture shape 찾기
        pics_to_replace = []
        for shape in list(slide.shapes):
            if shape.shape_type == 13:  # picture
                left_in = shape.left / 914400
                top_in = shape.top / 914400
                w_in = shape.width / 914400
                h_in = shape.height / 914400
                # SKC 로고 추정: 우측 상단 영역 (left > 9 인치, top < 1 인치) 또는 작은 사이즈
                if (left_in > 9 and top_in < 1 and w_in < 3) or \
                   (slide_idx == 1 and 0.5 < left_in < 4 and top_in > 6):
                    pics_to_replace.append((shape, left_in, top_in, w_in, h_in))

        for shape, left_in, top_in, w_in, h_in in pics_to_replace:
            # 위치/크기 보존하며 새 picture 추가 + 기존 삭제
            try:
                # 기존 SKC 로고 삭제
                sp = shape._element
                sp.getparent().remove(sp)
                # YSL 로고 같은 위치에 삽입
                slide.shapes.add_picture(
                    str(YSL_LOGO),
                    left=Emu(int(left_in * 914400)),
                    top=Emu(int(top_in * 914400)),
                    width=Emu(int(w_in * 914400)),
                    height=Emu(int(h_in * 914400)),
                )
                replaced += 1
            except Exception as e:
                print(f'슬라이드 {slide_idx} 로고 교체 실패: {e}')
    print(f'SKC 로고 → YSL 로고 교체: {replaced}개')


def main():
    print(f'SKC PPT 복사: {SKC_PPT.name} → {OUT_PPT.name}')
    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(SKC_PPT, OUT_PPT)

    prs = Presentation(OUT_PPT)
    print(f'총 슬라이드: {len(prs.slides)}')

    # 1. 모든 슬라이드 처리 (35슬라이드 그대로 처리 후 마지막에 잘라냄)
    for i, slide in enumerate(prs.slides, 1):
        try:
            process_slide(i, slide, prs)
            # 슬라이드 2 목차는 process_slide 후 다시 4섹션으로 갱신
            if i == 2:
                update_toc(slide)
            print(f'  슬라이드 {i} 처리 완료')
        except Exception as e:
            print(f'  슬라이드 {i} 처리 실패: {type(e).__name__}: {e}')

    # 2. 슬라이드 29~35 삭제 (5번 액션 섹션 + Appendix)
    print('\n슬라이드 29~35 삭제 중...')
    delete_slides(prs, list(range(29, 36)))
    print(f'삭제 후 총 슬라이드: {len(prs.slides)}')

    # 3. SKC 로고 → YSL 로고 교체
    print('\nSKC 로고 → YSL 로고 교체 중...')
    replace_skc_logo(prs)

    # 4. 모든 표 셀 병합 속성 일괄 제거 (vMerge로 행 가려지는 문제 방지)
    print('\n표 셀 병합 속성 제거 중...')
    unmerge_all_tables(prs)

    prs.save(OUT_PPT)
    print(f'\n저장 완료: {OUT_PPT}')


if __name__ == '__main__':
    main()
