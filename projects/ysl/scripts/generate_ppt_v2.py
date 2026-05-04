#!/usr/bin/env python3
"""YSL GEO Audit Report 생성 — V2 (CONSERVATIVE).

SKC 템플릿 구조 절대 건드리지 않고 텍스트/숫자만 교체.
- 표 셀 구조, 병합, 색상, 폰트 보존
- 빈 셀 클리어 안 함
- 컬럼 추가/삭제 안 함
- 단순히 cell.text 만 교체
"""
import json
import re
import shutil
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SKC_PPT = ROOT.parents[1] / 'archive' / 'BubbleShare_Skinceuticals_GEO Audit Report_1st_0416.pptx'
OUT_PPT = ROOT / 'final' / 'BubbleShare_YSL_GEO_Audit_Report_v2.pptx'

metrics = json.loads(Path('/tmp/ysl_metrics.json').read_text(encoding='utf-8'))
mention_df = pd.read_pickle('/tmp/ysl_mention.pkl')
citation_df = pd.read_pickle('/tmp/ysl_citation.pkl')


# =============================================================================
# 1. 텍스트 교체 (run-by-run, 다중 run 결합 처리)
# =============================================================================

REPLACEMENTS = [
    # 브랜드명
    ('SkinCeuticals', 'YSL Beauty'),
    ('skinceuticals', 'yslbeauty'),
    ('SKINCEUTICALS', 'YSL BEAUTY'),
    ('스킨수티컬즈', '입생로랑'),
    ('SKC', 'YSL'),
    # 시그니처 제품
    ('CE Ferulic', 'Libre EDP'),
    ('CE페룰릭', 'Libre'),
    ('CE 페룰릭', 'Libre'),
    # 카테고리 (SKC → YSL)
    ('항산화 앰플', '럭셔리 향수'),
    ('비타민C 앰플', '럭셔리 향수'),
]


def replace_in_paragraph(para, replacements):
    full = ''.join(r.text for r in para.runs)
    new = full
    for old, rep in replacements:
        new = new.replace(old, rep)
    if new != full and para.runs:
        para.runs[0].text = new
        for r in para.runs[1:]:
            r.text = ''


def replace_in_text_frame(tf, replacements):
    for para in tf.paragraphs:
        replace_in_paragraph(para, replacements)


def apply_text_replacements(slide, replacements):
    def walk(shape):
        if shape.has_text_frame:
            replace_in_text_frame(shape.text_frame, replacements)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    replace_in_text_frame(cell.text_frame, replacements)
        if shape.shape_type == 6:  # Group
            try:
                for inner in shape.shapes:
                    walk(inner)
            except Exception:
                pass
    for shape in slide.shapes:
        walk(shape)


# =============================================================================
# 2. 셀 값 안전 교체 (구조 보존)
# =============================================================================

def write_cell(cell, value):
    """셀 텍스트만 교체. 폰트/색/병합/테두리 보존."""
    tf = cell.text_frame
    # 첫 paragraph의 첫 run만 교체. 나머지 run은 빈문자열.
    if tf.paragraphs and tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].text = str(value)
        for r in tf.paragraphs[0].runs[1:]:
            r.text = ''
        # 추가 paragraph 비우기
        for para in tf.paragraphs[1:]:
            for run in para.runs:
                run.text = ''
    elif tf.paragraphs:
        run = tf.paragraphs[0].add_run()
        run.text = str(value)


# =============================================================================
# 3. 슬라이드별 데이터 채우기 (SKC 구조 그대로)
# =============================================================================

def fill_slide_4(slide):
    """슬라이드 4: Overview — 타임라인/플랫폼/쿼리 + Metric/Value 표."""
    tables = [s for s in slide.shapes if s.has_table]

    total_q = 360
    total_resp = metrics['total_responses']
    total_cit = metrics['total_citations']

    for tbl_shape in tables:
        t = tbl_shape.table
        rows, cols = len(t.rows), len(t.columns)

        # 2x3 표: 타임라인 / 플랫폼 / 쿼리
        if rows == 2 and cols == 3:
            write_cell(t.rows[1].cells[0], '분석: 2026-04-24 ~ 04-27\n수집 4사이클')
            write_cell(t.rows[1].cells[1], 'ChatGPT\nGoogle (AI Overview)\nNaver')
            write_cell(t.rows[1].cells[2], f'{total_q}건 질문\n+ {total_q}건 키워드')

        # 8x2 Metric/Value 표
        if rows == 8 and cols == 2 and 'Metric' in t.rows[0].cells[0].text_frame.text:
            metric_data = [
                ('Unique Queries', f'{total_q}'),
                ('Total Responses', f'{total_resp:,}'),
                ('Citation', f'{total_cit:,}'),
                ('Channels', 'ChatGPT, Google AIO, Naver'),
                ('Responses per Query (avg)', f'{total_resp/total_q:.1f}'),
                ('Analysis Period', '2026.04.24 ~ 2026.04.27'),
                ('Brand', 'YSL Beauty (입생로랑)'),
            ]
            for i, (k, v) in enumerate(metric_data):
                if i + 1 < rows:
                    write_cell(t.rows[i + 1].cells[0], k)
                    write_cell(t.rows[i + 1].cells[1], v)


def fill_slide_8(slide):
    """슬라이드 8: AI Engine mention rate 4개 표.

    SKC 원본 구조 (6x4 / 6x4 / 7x7 / 7x6):
    - 표 0 (Overall): 헤더 R0, 데이터 R1~
    - 표 1 (Channel by AI Engine): R0/R1 헤더, R2~ 데이터 (ChatGPT/Naver/Google/Total)
    - 표 2 (Category): R0/R1 헤더, R2~ 데이터
    - 표 3 (Intention): R0/R1 헤더, R2~ 데이터
    """
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) < 4:
        return

    f = metrics['funnel_all']
    total_cit = metrics['total_citations']
    ysl_cit = metrics['ysl_brand_url_count']

    # 표 0 (6x4 또는 7x4): Overall — Total/AI Existence/Brand Mention/Citation
    t0 = tables[0].table
    # SKC 원본: R0=title, R1=AI Existence(전체), R2=Commercial, R3=Brand Mention, R4=Citation
    # 데이터 R1부터 일관되게 채움
    overall_rows = [
        ('Total Query Set', f"{f['total']:,}", '100%', '쿼리 × AI Engine × 4 사이클'),
        ('AI Existence', f"{f['ai_existence']:,}", f"{f['ai_existence_rate']:.1f}%", '응답 생성된 쿼리'),
        ('Brand Mention', f"{f['brand_mention']:,}", f"{f['brand_mention_rate']:.1f}%", 'YSL 언급 응답'),
        ('Citation (전체)', f"{total_cit:,}", '—', '응답 내 인용 URL'),
        ('Citation (YSL 자사)', f"{ysl_cit}", f"{ysl_cit/max(total_cit,1)*100:.2f}%", 'YSL 자사몰 도메인'),
    ]
    for i, row in enumerate(overall_rows):
        r = i + 1
        if r >= len(t0.rows):
            break
        for c, val in enumerate(row):
            if c < len(t0.columns):
                write_cell(t0.rows[r].cells[c], val)

    # 표 1: Channel — SKC 원본 R0=AI Engine 헤더, R1=서브헤더(Mentions/Rate), R2~ 데이터
    # Total/Brand Mentions 컬럼 구조라 단순화: ChatGPT/Google AIO/Naver 3행
    t1 = tables[1].table
    chan_rows = []
    for r in metrics['by_channel']:
        ch = r['channel']
        label = {'chatgpt': 'ChatGPT', 'google': 'Google AIO', 'naver': 'Naver'}.get(ch, ch)
        chan_rows.append((label, f"{r['total']:,}", f"{r['mention']:,}", f"{r['rate']:.2f}%"))
    # SKC 원본은 데이터가 R2부터 (R0/R1은 헤더). 시도: R2부터 채움
    start_row = 2 if len(t1.rows) >= 5 else 1
    for i, row in enumerate(chan_rows):
        r = start_row + i
        if r >= len(t1.rows):
            break
        for c, val in enumerate(row):
            if c < len(t1.columns):
                write_cell(t1.rows[r].cells[c], val)

    # 표 2: Category Funnel
    t2 = tables[2].table
    cat_rows = []
    for r in metrics['by_category']:
        cat = r['category']
        cit = metrics['cit_by_category'].get(cat, 0)
        cit_share = cit / max(total_cit, 1) * 100
        cat_rows.append((cat, f"{r['total']:,}", f"{r['mention']:,}", f"{r['rate']:.2f}%",
                        f"{cit:,}", f"{cit_share:.1f}%", ''))
    start_row = 2 if len(t2.rows) >= 5 else 1
    for i, row in enumerate(cat_rows):
        r = start_row + i
        if r >= len(t2.rows):
            break
        for c, val in enumerate(row):
            if c < len(t2.columns):
                write_cell(t2.rows[r].cells[c], val)

    # 표 3: Intention Funnel
    t3 = tables[3].table
    int_rows = []
    for r in metrics['by_intent']:
        intent = r['intent']
        cit = metrics['cit_by_intent'].get(intent, 0)
        cit_per = cit / max(r['total'], 1)
        int_rows.append((intent, f"{r['total']:,}", f"{r['mention']:,}", f"{r['rate']:.2f}%",
                        f"{cit:,}", f"{cit_per:.1f}"))
    start_row = 2 if len(t3.rows) >= 5 else 1
    for i, row in enumerate(int_rows):
        r = start_row + i
        if r >= len(t3.rows):
            break
        for c, val in enumerate(row):
            if c < len(t3.columns):
                write_cell(t3.rows[r].cells[c], val)


def fill_slide_11(slide):
    """슬라이드 11: 채널 × 경쟁사 (SKC 원본 6x19, 데이터 R2~R5).

    SKC 구조:
    - R0+R1 헤더 (R0: AI Engine/Questions/Total Brand Mentions/SkinCeuticals/competitor1/2/...)
    - R0의 brand columns은 gridSpan=2 (Mentions + Rate 두 서브컬럼)
    - R1: sub-headers (Mentions / Rate %)
    - R2~R4: 데이터 행 (ChatGPT/Naver/Google)
    - R5: Grand Total

    YSL 데이터로 R2~R5 채움.
    """
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return
    t = tables[0].table

    sov_overall = metrics['sov_overall']
    sov_by_channel = metrics['sov_by_channel']

    # 8 brands (YSL + top 7 competitors by mention)
    brands = ['YSL Beauty', 'Dior', 'Chanel', 'Jo Malone', 'Hera', 'Sulwhasoo', 'MAC', 'Estee Lauder']

    # R0: brand 헤더 (col 3, 5, 7, 9, 11, 13, 15, 17 — gridSpan=2 구조 기준 짝수 위치)
    # SKC 원본은 col 3에 brand1, col 5에 brand2... 식으로 2칸씩 띄움
    # col 3=YSL, col 5=Dior, col 7=Chanel, ...
    for i, brand in enumerate(brands):
        col_brand = 3 + i * 2
        if col_brand < len(t.columns):
            write_cell(t.rows[0].cells[col_brand], brand)

    # 채널별 데이터 R2~R4
    chan_order = [('chatgpt', 'ChatGPT'), ('google', 'Google AIO'), ('naver', 'Naver')]
    for r_idx_offset, (chan_key, chan_label) in enumerate(chan_order):
        r_idx = 2 + r_idx_offset
        if r_idx >= len(t.rows):
            break
        chan_row = next((cr for cr in sov_by_channel if cr['channel'] == chan_key), None)
        if not chan_row:
            continue
        # C0: AI Engine 이름
        write_cell(t.rows[r_idx].cells[0], chan_label)
        # C1: Questions (응답 수)
        write_cell(t.rows[r_idx].cells[1], f"{chan_row['total']:,}")
        # C2: Total Brand Mentions (해당 채널 전체 브랜드 멘션 합)
        total_mentions = sum(int(chan_row.get(b, 0) * chan_row['total'] / 100) for b in brands)
        write_cell(t.rows[r_idx].cells[2], f"{total_mentions:,}")
        # C3+: 브랜드별 (Mentions, Rate) 짝
        for i, brand in enumerate(brands):
            col_m = 3 + i * 2
            col_r = col_m + 1
            rate = chan_row.get(brand, 0)
            count = int(rate * chan_row['total'] / 100)
            if col_m < len(t.columns):
                write_cell(t.rows[r_idx].cells[col_m], f"{count}")
            if col_r < len(t.columns):
                write_cell(t.rows[r_idx].cells[col_r], f"{rate:.2f}%")

    # R5: Grand Total
    if len(t.rows) >= 6:
        total_q = sum(cr['total'] for cr in sov_by_channel)
        write_cell(t.rows[5].cells[0], 'Grand Total')
        write_cell(t.rows[5].cells[1], f'{total_q:,}')
        total_brand_mentions = sum(int(sov_overall.get(b, 0) * total_q / 100) for b in brands)
        write_cell(t.rows[5].cells[2], f'{total_brand_mentions:,}')
        for i, brand in enumerate(brands):
            col_m = 3 + i * 2
            col_r = col_m + 1
            rate = sov_overall.get(brand, 0)
            count = int(rate * total_q / 100)
            if col_m < len(t.columns):
                write_cell(t.rows[5].cells[col_m], f"{count}")
            if col_r < len(t.columns):
                write_cell(t.rows[5].cells[col_r], f"{rate:.2f}%")


def fill_slide_19(slide):
    """슬라이드 19: SOV 랭킹 — 3개 11x5 표 (전체/카테고리A/카테고리B)."""
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) < 3:
        return

    sov = metrics['sov_overall']
    sov_by_cat = metrics['sov_by_category']

    def fill_sov_table(t, sov_dict):
        # SKC 원본은 R0+R1 헤더, R2~ 데이터 (Top 9 brands + Total)
        ranked = sorted(sov_dict.items(), key=lambda x: -x[1])
        for i, (brand, rate) in enumerate(ranked[:9]):
            r = i + 2
            if r >= len(t.rows):
                break
            write_cell(t.rows[r].cells[0], str(i + 1))
            write_cell(t.rows[r].cells[1], brand)
            write_cell(t.rows[r].cells[2], f"{rate:.2f}%")

    fill_sov_table(tables[0].table, sov)
    perfume = next((r for r in sov_by_cat if r.get('category') == '향수'), {})
    if perfume:
        fill_sov_table(tables[1].table,
                      {b: perfume.get(b, 0) for b in sov.keys()})
    gift = next((r for r in sov_by_cat if r.get('category') == '기프팅'), {})
    if gift:
        fill_sov_table(tables[2].table,
                      {b: gift.get(b, 0) for b in sov.keys()})


def fill_slide_24(slide):
    """슬라이드 24: 채널별 Top 7 도메인 — 8x5 표 3개."""
    tables = [s for s in slide.shapes if s.has_table]
    if len(tables) < 3:
        return

    plats = ['chatgpt', 'google', 'naver']
    for i, plat in enumerate(plats[:len(tables)]):
        if plat not in metrics['top_domains_by_platform']:
            continue
        t = tables[i].table
        domains = list(metrics['top_domains_by_platform'][plat].items())[:7]
        total = sum(metrics['top_domains_by_platform'][plat].values())
        for j, (d, c) in enumerate(domains):
            r = j + 2
            if r >= len(t.rows):
                break
            write_cell(t.rows[r].cells[0], str(j + 1))
            write_cell(t.rows[r].cells[1], d)
            write_cell(t.rows[r].cells[2], f"{c:,}")
            write_cell(t.rows[r].cells[3], f"{c/max(total,1)*100:.1f}%")


def fill_slide_26(slide):
    """슬라이드 26: 이커머스 인용 — 11x5 표."""
    tables = [s for s in slide.shapes if s.has_table]
    if not tables:
        return

    ecom = citation_df[citation_df['domain_type'] == 'Commerce']
    top = ecom['Domain'].value_counts().head(9)
    total_ecom = len(ecom)
    t = tables[0].table
    for j, (d, c) in enumerate(top.items()):
        r = j + 2
        if r >= len(t.rows):
            break
        write_cell(t.rows[r].cells[0], str(j + 1))
        write_cell(t.rows[r].cells[1], d)
        write_cell(t.rows[r].cells[2], f"{c}")
        write_cell(t.rows[r].cells[3], f"{c/max(total_ecom,1)*100:.1f}%")


# =============================================================================
# 4. 로고 교체
# =============================================================================

def replace_skc_logo(prs):
    YSL_LOGO = ROOT / 'assets' / 'ysl_logo_1.png'
    if not YSL_LOGO.exists():
        return
    replaced = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        pics = []
        for shape in list(slide.shapes):
            if shape.shape_type == 13:
                left_in = shape.left / 914400
                top_in = shape.top / 914400
                w_in = shape.width / 914400
                h_in = shape.height / 914400
                if (left_in > 9 and top_in < 1 and w_in < 3) or \
                   (slide_idx == 1 and 0.5 < left_in < 4 and top_in > 6):
                    pics.append((shape, left_in, top_in, w_in, h_in))
        for shape, left_in, top_in, w_in, h_in in pics:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(
                    str(YSL_LOGO),
                    left=Emu(int(left_in * 914400)),
                    top=Emu(int(top_in * 914400)),
                    width=Emu(int(w_in * 914400)),
                    height=Emu(int(h_in * 914400)),
                )
                replaced += 1
            except Exception:
                pass
    print(f'로고 교체: {replaced}개')


# =============================================================================
# 5. 슬라이드 삭제
# =============================================================================

def delete_slides(prs, indices_1based):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for idx in sorted(indices_1based, reverse=True):
        sld = slides[idx - 1]
        rId = sld.attrib['{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id']
        prs.part.drop_rel(rId)
        xml_slides.remove(sld)


# =============================================================================
# 메인
# =============================================================================

def main():
    print(f'SKC PPT 복사 → {OUT_PPT.name}')
    OUT_PPT.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy(SKC_PPT, OUT_PPT)
    prs = Presentation(OUT_PPT)
    print(f'총 슬라이드: {len(prs.slides)}')

    # 1. 모든 슬라이드 — 텍스트 일괄 교체
    print('\n[1] 텍스트 교체...')
    for slide in prs.slides:
        apply_text_replacements(slide, REPLACEMENTS)

    # 2. 슬라이드별 데이터 채움 (구조 보존)
    print('[2] 데이터 채움...')
    fill_slide_4(prs.slides[3])
    fill_slide_8(prs.slides[7])
    fill_slide_11(prs.slides[10])
    fill_slide_19(prs.slides[18])
    fill_slide_24(prs.slides[23])
    fill_slide_26(prs.slides[25])

    # 3. 슬라이드 삭제 (29~35: 5번 액션 + Appendix)
    print('[3] 슬라이드 29~35 삭제...')
    delete_slides(prs, list(range(29, 36)))
    print(f'   → 총 슬라이드: {len(prs.slides)}')

    # 4. 로고 교체
    print('[4] 로고 교체...')
    replace_skc_logo(prs)

    prs.save(OUT_PPT)
    print(f'\n저장: {OUT_PPT}')


if __name__ == '__main__':
    main()
